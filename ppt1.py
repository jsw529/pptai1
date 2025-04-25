import streamlit as st
from pptx import Presentation
from google.cloud import texttospeech
from moviepy.editor import *
import tempfile
import os

# TTS 생성 함수
def generate_speech(text, language_code='ko-KR', voice_gender='NEUTRAL', voice_name='ko-KR-Wavenet-A'):
    client = texttospeech.TextToSpeechClient()

    input_text = texttospeech.SynthesisInput(text=text)
    voice = texttospeech.VoiceSelectionParams(
        language_code=language_code,
        ssml_gender=texttospeech.SsmlVoiceGender[voice_gender],
        name=voice_name
    )
    audio_config = texttospeech.AudioConfig(audio_encoding=texttospeech.AudioEncoding.MP3)

    response = client.synthesize_speech(input=input_text, voice=voice, audio_config=audio_config)

    with tempfile.NamedTemporaryFile(delete=False) as tmp_audio:
        tmp_audio.write(response.audio_content)
        return tmp_audio.name

# PPT 텍스트 추출 함수
def extract_text_from_ppt(ppt_file):
    ppt = Presentation(ppt_file)
    text_content = []
    for slide in ppt.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_content.append(slide_text)
    return text_content

# 영상 생성 함수
def create_video_from_audio(audio_file, slides_text, output_path="output_video.mp4"):
    try:
        clips = []
        for slide_text in slides_text:
            # 각 슬라이드마다 5초씩 비디오 클립을 생성
            clip = ColorClip(size=(1280, 720), color=(255, 255, 255), duration=5)
            txt_clip = TextClip(slide_text, fontsize=50, color='black', size=(1280, 720), method='caption')
            txt_clip = txt_clip.set_position('center').set_duration(5)
            clips.append(clip.set_duration(5).fx(lambda clip: clip.set_audio(AudioFileClip(audio_file))))

        final_video = concatenate_videoclips(clips)
        
        # 영상 파일 저장
        final_video.write_videofile(output_path, codec="libx264")
        
        return output_path  # 영상 파일 경로 반환
    except Exception as e:
        st.error(f"영상 생성 중 오류 발생: {e}")
        return None

# Streamlit 앱 인터페이스
st.title('PPT 음성 합성 및 영상 생성기')

# PPT 파일 업로드
ppt_file = st.file_uploader("PPT 파일을 업로드하세요", type=["pptx"])

# 음성 톤 선택 (여러 톤 중 선택)
voice_gender = st.selectbox(
    "음성 톤을 선택하세요:",
    ("NEUTRAL", "MALE", "FEMALE")
)

voice_name = ""
if voice_gender == "NEUTRAL":
    voice_name = "ko-KR-Wavenet-A"
elif voice_gender == "MALE":
    voice_name = "ko-KR-Wavenet-B"
elif voice_gender == "FEMALE":
    voice_name = "ko-KR-Wavenet-C"

if ppt_file:
    # PPT에서 텍스트 추출
    text_content = extract_text_from_ppt(ppt_file)
    st.write("슬라이드 내용:")
    for idx, text in enumerate(text_content):
        st.text_area(f"슬라이드 {idx + 1}", value=text, height=100)

    # 대본을 AI 음성으로 변환 버튼
    if st.button('음성 생성'):
        audio_file = generate_speech("\n".join(text_content), voice_gender=voice_gender, voice_name=voice_name)
        st.audio(audio_file)

        # 영상 생성 버튼
        if st.button('영상 생성'):
            video_file = "output_video.mp4"
            st.text("영상 생성 중...")

            # 영상 생성 처리 시간 표시
            video_file_path = create_video_from_audio(audio_file, text_content, output_path=video_file)
            
            if video_file_path:
                st.video(video_file_path)
                st.success('영상 생성 완료!')

                # 영상 다운로드 버튼 추가
                with open(video_file_path, "rb") as video_file_data:
                    st.download_button(
                        label="영상 다운로드",
                        data=video_file_data,
                        file_name="output_video.mp4",
                        mime="video/mp4"
                    )
            else:
                st.error('영상 생성에 실패했습니다.')
